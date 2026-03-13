"""
PPT 生成服务
使用 python-pptx 根据论文内容生成 PowerPoint 演示文稿
"""

import os
import re
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTService:
    """PPT 生成服务"""

    def __init__(self, output_dir: str = "./temp/ppt_output"):
        """
        初始化 PPT 服务

        Args:
            output_dir: PPT 输出目录
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def parse_paper_structure(self, markdown_content: str) -> Dict:
        """
        解析论文的 Markdown 内容,提取结构化信息

        Args:
            markdown_content: 论文的 Markdown 内容

        Returns:
            包含标题、作者、章节等信息的字典
        """
        structure = {
            'title': '未命名论文',
            'authors': [],
            'sections': [],
            'images': []
        }

        lines = markdown_content.split('\n')
        current_section = None
        current_content = []

        for line in lines:
            # 提取一级标题（论文标题）
            if line.startswith('# ') and not structure['title'] or structure['title'] == '未命名论文':
                structure['title'] = line[2:].strip()
                continue

            # 提取作者信息
            if re.match(r'^[A-Z][a-z]+ [A-Z]', line) and 'University' in line:
                structure['authors'].append(line.strip())
                continue

            # 提取章节（二级标题）
            if line.startswith('## '):
                # 保存上一章节
                if current_section:
                    structure['sections'].append({
                        'title': current_section,
                        'content': '\n'.join(current_content).strip()
                    })

                # 开始新章节
                current_section = line[3:].strip()
                current_content = []
            else:
                current_content.append(line)

        # 保存最后一个章节
        if current_section:
            structure['sections'].append({
                'title': current_section,
                'content': '\n'.join(current_content).strip()
            })

        # 提取图片
        pattern = r'!\[\]\(images/([^)]+)\)'
        matches = re.findall(pattern, markdown_content)
        structure['images'] = matches

        return structure

    def create_title_slide(self, prs: Presentation, title: str, authors: List[str]):
        """
        创建标题页

        Args:
            prs: Presentation 对象
            title: 论文标题
            authors: 作者列表
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 手动添加标题文本框
        title_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(2.5),
            width=Inches(8),
            height=Inches(1)
        )

        # 设置标题文本和格式
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 78, 120)
        title_para.alignment = PP_ALIGN.CENTER

        # 添加作者信息
        if authors:
            author_box = slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(4),
                width=Inches(8),
                height=Inches(1)
            )
            author_frame = author_box.text_frame
            author_frame.text = '\n'.join(authors[:5])  # 最多显示5个作者

            for paragraph in author_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(89, 89, 89)
                paragraph.alignment = PP_ALIGN.CENTER

    def create_outline_slide(self, prs: Presentation, sections: List[Dict]):
        """
        创建目录页

        Args:
            prs: Presentation 对象
            sections: 章节列表
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 手动添加标题
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(9),
            height=Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "目录"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 78, 120)

        # 添加目录内容
        content_box = slide.shapes.add_textbox(
            left=Inches(1),
            top=Inches(2),
            width=Inches(8),
            height=Inches(5)
        )
        content_frame = content_box.text_frame

        for i, section in enumerate(sections[:8]):  # 最多显示8个章节
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = f"{i + 1}. {section['title']}"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(68, 114, 196)
            p.space_before = Pt(12)
            p.space_after = Pt(6)

    def create_section_slide(self, prs: Presentation, section: Dict, section_num: int, images_dir: Optional[Path] = None):
        """
        创建章节内容页

        Args:
            prs: Presentation 对象
            section: 章节信息
            section_num: 章节编号
            images_dir: 图片目录（可选）
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 手动添加章节标题
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(9),
            height=Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"{section_num}. {section['title']}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 78, 120)

        # 处理内容：提取要点和图片
        content_lines = section['content'].split('\n')
        bullet_points = []
        image_path = None

        for line in content_lines:
            line = line.strip()
            # 跳过空行和标题行
            if not line or line.startswith('#'):
                continue

            # 提取图片
            if line.startswith('![](images/') or '![](images\\' in line:
                import re
                img_match = re.search(r'!\[\]\(images[\\/]([^)]+)\)', line)
                if img_match and images_dir:
                    image_path = images_dir / img_match.group(1)
                continue

            # 提取列表项和重要段落
            if line.startswith('-') or line.startswith('*') or len(line) > 20:
                # 清理 markdown 标记
                clean_line = line.lstrip('-*').strip()
                if clean_line:
                    bullet_points.append(clean_line)

        # 布局策略：左文字右图片
        has_image = image_path is not None and image_path.exists()
        has_points = len(bullet_points) > 0

        if has_image and has_points:
            # 左右布局：文字左边，图片右边

            # 先添加文字框（在左侧）
            content_box = slide.shapes.add_textbox(
                left=Inches(0.5),      # 左边距 0.5 英寸
                top=Inches(1.5),       # 上边距 1.5 英寸
                width=Inches(4.5),     # 宽度 4.5 英寸（左侧占45%）
                height=Inches(4.0)     # 高度 4 英寸
            )

            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # 添加要点（每页最多6个要点）
            for i, point in enumerate(bullet_points[:6]):
                if i > 0:
                    content_frame.add_paragraph()

                p = content_frame.paragraphs[i]
                p.text = f"• {point}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(89, 89, 89)
                p.space_before = Pt(8)
                p.space_after = Pt(6)

            # 后添加图片（在右侧），这样图片会在上层
            try:
                from PIL import Image

                with Image.open(str(image_path)) as img:
                    img_width, img_height = img.size
                    img_aspect_ratio = img_height / img_width

                # 图片在右侧，占45%宽度
                img_max_width = 4.3
                img_max_height = 4.0

                if img_aspect_ratio > img_max_height / img_max_width:
                    img_height_display = img_max_height
                    img_width_display = img_height_display / img_aspect_ratio
                else:
                    img_width_display = img_max_width
                    img_height_display = img_width_display * img_aspect_ratio

                # 图片位置：左侧从 5.2 英寸开始（文字框 0.5 + 4.5 = 5.0，留 0.2 间隙）
                img_left = 5.2
                img_top = 1.5 + (4.0 - img_height_display) / 2

                # 添加图片
                pic = slide.shapes.add_picture(
                    str(image_path),
                    left=Inches(img_left),
                    top=Inches(img_top),
                    width=Inches(img_width_display)
                )

            except Exception as e:
                print(f"添加图片失败: {e}")

        elif has_image and not has_points:
            # 只有图片：居中显示
            try:
                from PIL import Image

                with Image.open(str(image_path)) as img:
                    img_width, img_height = img.size
                    img_aspect_ratio = img_height / img_width

                max_img_height = 4.5
                max_img_width = 9

                if img_aspect_ratio > max_img_height / max_img_width:
                    img_height_display = max_img_height
                    img_width_display = img_height_display / img_aspect_ratio
                else:
                    img_width_display = max_img_width
                    img_height_display = img_width_display * img_aspect_ratio

                img_left = (10 - img_width_display) / 2
                img_top = 1.5 + (4.5 - img_height_display) / 2

                pic = slide.shapes.add_picture(
                    str(image_path),
                    left=Inches(img_left),
                    top=Inches(img_top),
                    width=Inches(img_width_display)
                )

                pic.top = Inches(img_top)
                pic.left = Inches(img_left)

            except Exception as e:
                print(f"添加图片失败: {e}")

        elif has_points and not has_image:
            # 只有文字：占满整个宽度
            content_box = slide.shapes.add_textbox(
                left=Inches(0.7),
                top=Inches(1.5),
                width=Inches(8.6),
                height=Inches(5)
            )

            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # 添加要点（每页最多6个要点）
            for i, point in enumerate(bullet_points[:6]):
                if i > 0:
                    content_frame.add_paragraph()

                p = content_frame.paragraphs[i]
                p.text = point
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(89, 89, 89)
                p.space_before = Pt(8)
                p.space_after = Pt(6)

    def create_summary_slide(self, prs: Presentation, paper_structure: Dict):
        """
        创建总结页

        Args:
            prs: Presentation 对象
            paper_structure: 论文结构信息
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 手动添加标题
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.5),
            width=Inches(9),
            height=Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "总结"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(31, 78, 120)

        # 添加总结内容
        content_box = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=Inches(1.5),
            width=Inches(8.6),
            height=Inches(5)
        )
        content_frame = content_box.text_frame

        summary_points = [
            f"**论文标题**: {paper_structure['title']}",
            f"**主要贡献**: 见后续各章节详述",
            f"**核心方法**: 详见方法章节",
            f"**实验结果**: 详见实验章节",
            "",
            "**下一步工作**:",
            "- 深入理解论文细节",
            "- 复现实验结果",
            "- 思考应用场景"
        ]

        for i, point in enumerate(summary_points):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            p.text = point
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(89, 89, 89)
            p.space_before = Pt(10)
            p.space_after = Pt(8)

    async def generate_ppt(self, markdown_content: str, task_id: str, images_dir: Optional[Path] = None) -> Dict:
        """
        生成 PPT

        Args:
            markdown_content: 论文的 Markdown 内容
            task_id: 任务 ID
            images_dir: 图片目录（可选）

        Returns:
            生成结果和文件路径
        """
        try:
            # 解析论文结构
            paper_structure = self.parse_paper_structure(markdown_content)

            # 创建 Presentation 对象
            prs = Presentation()

            # 设置幻灯片大小（16:9）
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            # 1. 创建标题页
            self.create_title_slide(prs, paper_structure['title'], paper_structure['authors'])

            # 2. 创建目录页
            if paper_structure['sections']:
                self.create_outline_slide(prs, paper_structure['sections'])

            # 3. 创建章节内容页
            for i, section in enumerate(paper_structure['sections'][:10], 1):  # 最多10个章节
                self.create_section_slide(prs, section, i, images_dir)

            # 4. 创建总结页
            self.create_summary_slide(prs, paper_structure)

            # 保存 PPT 文件
            output_path = self.output_dir / f"{task_id}.pptx"
            prs.save(str(output_path))

            return {
                "status": "success",
                "ppt_path": str(output_path),
                "slides_count": len(prs.slides),
                "title": paper_structure['title'],
                "sections_count": len(paper_structure['sections'])
            }

        except Exception as e:
            import traceback
            print(f"生成 PPT 时发生异常:\n{traceback.format_exc()}")
            return {
                "status": "error",
                "error": f"生成 PPT 失败: {str(e)}"
            }

    def parse_ppt_content(self, ppt_content: str) -> Dict:
        """
        解析 LLM 生成的 PPT 专用内容

        Args:
            ppt_content: LLM 生成的 PPT 内容（Markdown 格式）

        Returns:
            包含标题、作者、幻灯片列表等信息的字典
        """
        structure = {
            'title': '未命名演示文稿',
            'authors': [],
            'slides': []
        }

        lines = ppt_content.split('\n')
        current_slide = None
        current_title = None
        current_points = []
        current_image = None

        for line in lines:
            line = line.strip()

            # 提取主标题（一级标题）
            if line.startswith('# ') and not line.startswith('## '):
                structure['title'] = line[2:].strip()
                continue

            # 提取作者信息
            if line.startswith('## 作者信息'):
                # 下一行应该是作者列表
                continue
            elif line.startswith('## 目录'):
                # 跳过目录部分
                current_slide = None
                continue

            # 提取幻灯片（二级标题）
            if line.startswith('## 第') and '页' in line:
                # 保存上一张幻灯片
                if current_slide and current_title:
                    structure['slides'].append({
                        'title': current_title,
                        'points': current_points,
                        'image': current_image
                    })

                # 开始新幻灯片
                current_slide = line[3:].strip()
                current_title = None
                current_points = []
                current_image = None

            # 提取幻灯片标题
            elif line.startswith('**标题**') or line.startswith('**标题：') or line.startswith('**标题：**'):
                current_title = line.split('：')[-1].split(':')[-1].strip().strip('*').strip()

            # 提取要点
            elif line.startswith('- '):
                point = line[2:].strip()
                if point:
                    current_points.append(point)

            # 提取图片
            elif line.startswith('**图片**') or '![](images/' in line:
                import re
                img_match = re.search(r'!\[\]\(images/([^)]+)\)', line)
                if img_match:
                    current_image = img_match.group(1)

        # 保存最后一张幻灯片
        if current_title or current_points:
            structure['slides'].append({
                'title': current_title or current_slide,
                'points': current_points,
                'image': current_image
            })

        return structure

    def create_ppt_from_content(self, ppt_structure: Dict, task_id: str, images_dir: Optional[Path] = None) -> Dict:
        """
        基于解析后的 PPT 结构创建 PowerPoint 文件

        Args:
            ppt_structure: 解析后的 PPT 结构
            task_id: 任务 ID
            images_dir: 图片目录（可选）

        Returns:
            生成结果和文件路径
        """
        try:
            # 创建 Presentation 对象
            prs = Presentation()

            # 设置幻灯片大小（16:9）
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            # 1. 创建标题页
            title_slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 添加标题文本框
            title_box = title_slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(2),
                width=Inches(8),
                height=Inches(1.2)
            )

            title_frame = title_box.text_frame
            title_frame.text = ppt_structure['title']
            title_frame.word_wrap = True

            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(31, 78, 120)
            title_para.alignment = PP_ALIGN.CENTER

            # 2. 创建内容页
            for slide_data in ppt_structure['slides']:
                if not slide_data['title'] and not slide_data['points']:
                    continue

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                # 判断是否有图片
                has_image = bool(slide_data['image']) and images_dir
                has_points = bool(slide_data['points'])

                # DEBUG
                if slide_data['image']:
                    print(f"[DEBUG] 幻灯片'{slide_data['title'][:20]}': has_image={has_image}, image={slide_data['image'][:20]}, images_dir={images_dir is not None}")

                # 添加标题（始终在顶部）
                if slide_data['title']:
                    title_box = slide.shapes.add_textbox(
                        left=Inches(0.5),
                        top=Inches(0.5),
                        width=Inches(9),
                        height=Inches(0.8)
                    )

                    title_frame = title_box.text_frame
                    title_frame.text = slide_data['title']
                    title_frame.word_wrap = True

                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(32)
                    title_para.font.bold = True
                    title_para.font.color.rgb = RGBColor(31, 78, 120)

                # 布局策略：文字在左，图片在右
                if has_image and has_points:
                    # 左右布局：文字左边，图片右边
                    image_path = images_dir / slide_data['image']

                    # 先添加文字框（在左侧）
                    content_box = slide.shapes.add_textbox(
                        left=Inches(0.5),      # 左边距 0.5 英寸
                        top=Inches(1.5),       # 上边距 1.5 英寸
                        width=Inches(4.5),     # 宽度 4.5 英寸（左侧占45%）
                        height=Inches(4.0)     # 高度 4 英寸
                    )

                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True

                    for i, point in enumerate(slide_data['points']):
                        if i > 0:
                            content_frame.add_paragraph()

                        p = content_frame.paragraphs[i]
                        p.text = f"• {point}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(89, 89, 89)
                        p.space_before = Pt(8)
                        p.space_after = Pt(6)

                    # 后添加图片（在右侧），这样图片会在上层
                    if image_path.exists():
                        try:
                            from PIL import Image

                            with Image.open(str(image_path)) as img:
                                img_width, img_height = img.size
                                img_aspect_ratio = img_height / img_width

                            # 图片在右侧，占45%宽度
                            img_max_width = 4.3
                            img_max_height = 4.0

                            if img_aspect_ratio > img_max_height / img_max_width:
                                img_height_display = img_max_height
                                img_width_display = img_height_display / img_aspect_ratio
                            else:
                                img_width_display = img_max_width
                                img_height_display = img_width_display * img_aspect_ratio

                            # 图片位置：左侧从 5.2 英寸开始（文字框 0.5 + 4.5 = 5.0，留 0.2 间隙）
                            img_left = 5.2
                            img_top = 1.5 + (4.0 - img_height_display) / 2

                            # 添加图片
                            pic = slide.shapes.add_picture(
                                str(image_path),
                                left=Inches(img_left),
                                top=Inches(img_top),
                                width=Inches(img_width_display)
                            )

                            print(f"[IMAGE] '{slide_data['title'][:20]}': left={img_left:.2f}\", top={img_top:.2f}\", width={img_width_display:.2f}\", height={img_height_display:.2f}\"")

                        except Exception as e:
                            print(f"添加图片失败 {slide_data['image']}: {e}")
                            import traceback
                            traceback.print_exc()

                elif has_image and not has_points:
                    # 只有图片：居中显示
                    image_path = images_dir / slide_data['image']
                    if image_path.exists():
                        try:
                            from PIL import Image

                            with Image.open(str(image_path)) as img:
                                img_width, img_height = img.size
                                img_aspect_ratio = img_height / img_width

                            max_img_height = 4.5
                            max_img_width = 9

                            if img_aspect_ratio > max_img_height / max_img_width:
                                img_height_display = max_img_height
                                img_width_display = img_height_display / img_aspect_ratio
                            else:
                                img_width_display = max_img_width
                                img_height_display = img_width_display * img_aspect_ratio

                            img_left = (10 - img_width_display) / 2
                            img_top = 1.5 + (4.5 - img_height_display) / 2

                            pic = slide.shapes.add_picture(
                                str(image_path),
                                left=Inches(img_left),
                                top=Inches(img_top),
                                width=Inches(img_width_display)
                            )

                            pic.top = Inches(img_top)
                            pic.left = Inches(img_left)

                            print(f"[IMAGE] '{slide_data['title'][:20]}': pos=({img_left:.2f}\", {img_top:.2f}\"), size={img_width_display:.2f}\"x{img_height_display:.2f}\"")

                        except Exception as e:
                            print(f"添加图片失败 {slide_data['image']}: {e}")
                            import traceback
                            traceback.print_exc()

                elif has_points and not has_image:
                    # 只有文字：占满整个宽度
                    content_box = slide.shapes.add_textbox(
                        left=Inches(0.7),
                        top=Inches(1.5),
                        width=Inches(8.6),
                        height=Inches(4)
                    )

                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True

                    for i, point in enumerate(slide_data['points']):
                        if i > 0:
                            content_frame.add_paragraph()

                        p = content_frame.paragraphs[i]
                        p.text = f"• {point}"
                        p.font.size = Pt(20)
                        p.font.color.rgb = RGBColor(89, 89, 89)
                        p.space_before = Pt(10)
                        p.space_after = Pt(6)

            # 保存 PPT 文件
            output_path = self.output_dir / f"{task_id}.pptx"
            prs.save(str(output_path))

            return {
                "status": "success",
                "ppt_path": str(output_path),
                "slides_count": len(prs.slides),
                "title": ppt_structure['title']
            }

        except Exception as e:
            import traceback
            print(f"创建 PPT 时发生异常:\n{traceback.format_exc()}")
            return {
                "status": "error",
                "error": f"创建 PPT 失败: {str(e)}"
            }

    async def generate_ppt_from_notes(self, notes_content: str, task_id: str, images_dir: Optional[Path] = None) -> Dict:
        """
        基于笔记内容生成 PPT（新方法）

        Args:
            notes_content: LLM 生成的笔记内容
            task_id: 任务 ID
            images_dir: 图片目录（可选）

        Returns:
            生成结果和文件路径
        """
        try:
            # 解析笔记内容，提取 PPT 结构
            ppt_structure = self.parse_ppt_content(notes_content)

            # 创建 PPT
            result = self.create_ppt_from_content(ppt_structure, task_id, images_dir)

            return result

        except Exception as e:
            import traceback
            print(f"基于笔记生成 PPT 时发生异常:\n{traceback.format_exc()}")
            return {
                "status": "error",
                "error": f"生成 PPT 失败: {str(e)}"
            }

    def cleanup_task_files(self, task_id: str):
        """
        清理任务相关的 PPT 文件

        Args:
            task_id: 任务 ID
        """
        try:
            ppt_file = self.output_dir / f"{task_id}.pptx"
            if ppt_file.exists():
                ppt_file.unlink()
        except Exception as e:
            print(f"清理 PPT 文件失败: {e}")
