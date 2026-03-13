"""
测试PPT左文字右图片布局功能（带真实图片）
"""

import asyncio
from pathlib import Path
from services.ppt_service import PPTService

# 测试用的笔记内容（包含图片）
TEST_NOTES = '''
# 深度学习模型优化研究

## 第1页：研究背景

**标题**：研究背景

- 深度学习在各领域广泛应用
- 模型规模持续增长
- 推理效率成为关键问题

**图片**：![](images/test1.png)

## 第2页：纯文字页面

**标题**：现有方法

- 量化方法：降低数值精度
- 剪枝技术：移除冗余连接
- 知识蒸馏：Teacher-Student框架
- 动态推理：根据输入调整计算

## 第3页：方法概述

**标题**：我们的方法

- 自适应量化策略
- 动态剪枝算法
- 端到端优化框架

**图片**：![](images/test2.png)

## 第4页：结论

**标题**：总结

- 实现了2.3倍加速
- 精度损失仅0.13%
- 适用多种模型架构
'''


async def main():
    print("=" * 60)
    print("Testing PPT Layout: Left Text, Right Image")
    print("=" * 60)

    # Create PPT service
    ppt_service = PPTService()

    # Parse notes
    print("\n[1] Parsing notes content...")
    structure = ppt_service.parse_ppt_content(TEST_NOTES)

    print(f"   Title: {structure['title']}")
    print(f"   Total slides: {len(structure['slides'])}")

    for i, slide in enumerate(structure['slides'], 1):
        img_status = "[IMG]" if slide['image'] else "[TEXT]"
        print(f"   Slide {i}: {slide['title'][:25]} | Points: {len(slide['points'])} | {img_status}")

    # Generate PPT with images
    print("\n[2] Generating PPT with images...")
    task_id = "test_layout_with_images"
    images_dir = Path("test_images")  # Use our test images directory

    result = await ppt_service.generate_ppt_from_notes(
        notes_content=TEST_NOTES,
        task_id=task_id,
        images_dir=images_dir
    )

    if result["status"] == "success":
        print(f"   [OK] PPT generated successfully!")
        print(f"   File: {result['ppt_path']}")
        print(f"   Slides: {result['slides_count']}")

        # Verify file
        ppt_file = Path(result['ppt_path'])
        if ppt_file.exists():
            file_size = ppt_file.stat().st_size
            print(f"   Size: {file_size:,} bytes")

        # Check layout of specific slides
        print("\n[3] Verifying layout...")
        from pptx import Presentation
        prs = Presentation(ppt_file)

        # Slide 1 (title + text + image)
        if len(prs.slides) > 0:
            slide = prs.slides[0]
            print(f"\n   Slide 1 (Title): {len(slide.shapes)} shapes")

        # Slide 2 (should be: text left, image right)
        if len(prs.slides) > 1:
            slide = prs.slides[1]
            print(f"\n   Slide 2 (Text + Image): {len(slide.shapes)} shapes")
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text'):
                    has_text = "Yes" if shape.text.strip() else "No"
                else:
                    has_text = "No"
                print(f"      Shape {j}: pos=({shape.left.inches:.1f}\", {shape.top.inches:.1f}\") "
                      f"size={shape.width.inches:.1f}\"x{shape.height.inches:.1f}\" has_text={has_text}")

                # Check if this is the text box (should be on left)
                if hasattr(shape, 'text') and shape.text:
                    if shape.left.inches < 3.0:
                        print(f"        -> This is LEFT text box (correct!)")
                    else:
                        print(f"        -> This is RIGHT text box (might be image)")

        # Slide 3 (text only - should be full width)
        if len(prs.slides) > 2:
            slide = prs.slides[2]
            print(f"\n   Slide 3 (Text only): {len(slide.shapes)} shapes")
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text:
                    print(f"      Text box {j}: width={shape.width.inches:.1f}\" (should be ~8.6 for full width)")

    else:
        print(f"   [ERROR] {result.get('error')}")

    print("\n" + "=" * 60)
    print("Test completed!")
    print("=" * 60)


if __name__ == "__main__":
    asyncio.run(main())
